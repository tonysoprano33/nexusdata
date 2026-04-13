"""
Generador de reportes PowerPoint para análisis de datos
"""

import os
from datetime import datetime
from typing import Dict, Any, List
import io

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx no instalado. PowerPoint no disponible.")

def generate_pptx_report(analysis_data: Dict[str, Any]) -> bytes:
    """
    Genera un reporte PowerPoint completo del análisis
    """
    if not PPTX_AVAILABLE:
        return None
    
    prs = Presentation()
    
    # Colores tema
    COLOR_PRIMARY = RGBColor(99, 102, 241)  # Indigo
    COLOR_TEXT = RGBColor(255, 255, 255)    # Blanco
    COLOR_BG = RGBColor(23, 23, 23)         # Dark bg
    
    # Slide 1: Título
    title_slide_layout = prs.slide_layouts[0]  # Título
    slide = prs.slides.add_slide(title_slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "📊 Reporte de Análisis de Datos"
    subtitle.text = f"Generado por NexusData AI\n{datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # Estilo título
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    
    # Slide 2: Resumen Ejecutivo
    content_slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(content_slide_layout)
    
    slide.shapes.title.text = "📈 Resumen Ejecutivo"
    
    result = analysis_data.get('result', {})
    summary = result.get('summary', {})
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Métricas Principales:"
    
    metrics = [
        f"• Filas analizadas: {summary.get('total_rows', 0):,}",
        f"• Columnas: {summary.get('total_columns', 0)}",
        f"• Calidad de datos: {summary.get('data_quality_score', 0)}%",
        f"• Anomalías detectadas: {result.get('anomaly_detection', {}).get('detected_rows', 0)}",
    ]
    
    for metric in metrics:
        p = tf.add_paragraph()
        p.text = metric
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 3: Insights de IA
    insights = result.get('business_insights', '')
    if insights and 'Insight Engine Disabled' not in insights:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "🧠 Insights de Inteligencia Artificial"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Limpiar markdown
        insights_clean = insights.replace('**', '').replace('#', '').replace('*', '').replace('📊', '').replace('💡', '').replace('⚠️', '')
        
        # Dividir en líneas y agregar
        lines = [line.strip() for line in insights_clean.split('\n') if line.strip()]
        if lines:
            tf.text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
    
    # Slide 4: Gráficos Recomendados
    charts = result.get('chart_recommendations', [])
    if charts:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "📊 Visualizaciones Recomendadas"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Gráficos generados automáticamente:"
        
        for i, chart in enumerate(charts[:5], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {chart.get('title', 'Gráfico')} ({chart.get('type', 'chart')})"
            p.level = 1
            p.font.size = Pt(16)
            
            # Agregar insight si existe
            insight = chart.get('insight', '')
            if insight:
                p2 = tf.add_paragraph()
                p2.text = f"   → {insight[:60]}..."
                p2.level = 2
                p2.font.size = Pt(12)
                p2.font.italic = True
    
    # Slide 5: Estadísticas (muestra)
    stats = result.get('descriptive_statistics', {})
    if stats:
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = "📉 Estadísticas Descriptivas"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "Análisis estadístico de variables numéricas:"
        
        for col_name, col_stats in list(stats.items())[:3]:
            p = tf.add_paragraph()
            p.text = f"• {col_name}:"
            p.level = 1
            p.font.size = Pt(16)
            p.font.bold = True
            
            if isinstance(col_stats, dict):
                stats_text = ', '.join([f"{k}: {v:.2f}" if isinstance(v, (int, float)) else f"{k}: {v}" 
                                       for k, v in list(col_stats.items())[:4]])
                p2 = tf.add_paragraph()
                p2.text = f"  {stats_text}"
                p2.level = 2
                p2.font.size = Pt(12)
    
    # Slide 6: Footer/Contacto
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = "✨ Análisis Completo"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Este reporte fue generado automáticamente por NexusData AI"
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "🚀 Análisis avanzado sin costos ni límites"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Visita: https://nexusdata.ai"
    p.font.size = Pt(14)
    
    # Guardar a bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    pptx_bytes = buffer.getvalue()
    buffer.close()
    
    return pptx_bytes
