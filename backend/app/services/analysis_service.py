from __future__ import annotations

import io
import logging
import math
import traceback
import uuid
from typing import Any, Dict, List, Optional, Tuple

import numpy as np
import pandas as pd
from pandas.api.types import is_datetime64_any_dtype, is_numeric_dtype, is_string_dtype

from app.db.supabase import get_supabase_db
from app.services.data_cleaning_service import get_data_cleaning_service
from app.services.gemini_service import GeminiService
from app.services.groq_service import GroqService

try:  # pragma: no cover - optional ML extras
    from sklearn.cluster import KMeans
    from sklearn.ensemble import IsolationForest
    from sklearn.linear_model import LinearRegression
    from sklearn.metrics import r2_score
    from sklearn.model_selection import train_test_split
    from sklearn.preprocessing import StandardScaler

    SKLEARN_AVAILABLE = True
except Exception:  # pragma: no cover - optional dependency
    SKLEARN_AVAILABLE = False

logger = logging.getLogger(__name__)


class AnalysisService:
    """Orchestrates parsing, cleaning, analysis, storage, exports, and chat."""

    def __init__(self):
        self.gemini = GeminiService()
        self.groq = GroqService()
        self.cleaning_service = get_data_cleaning_service()

    async def upload_and_analyze(
        self,
        file_content: bytes,
        filename: str,
        provider: str = "gemini",
        custom_prompt: Optional[str] = None,
        preview: Optional[list] = None,
    ) -> Dict[str, Any]:
        """Upload a dataset, clean it, analyze it, and persist the result."""
        del preview  # legacy argument kept for compatibility

        db = get_supabase_db()
        analysis_id = str(uuid.uuid4())
        await db.create_analysis(analysis_id, filename)

        try:
            df_raw = self._parse_dataset(file_content, filename)
            df_clean, cleaning_report_raw = self.cleaning_service.clean_dataset(df_raw)
            cleaning_report = self._normalize_cleaning_report(cleaning_report_raw)
            await db.update_analysis_status(analysis_id, "processing")

            raw_stats = self._build_dataset_statistics(df_raw)
            clean_stats = self._build_dataset_statistics(df_clean)
            anomaly_detection = self._detect_anomalies(df_clean)
            advanced_analytics = self._generate_advanced_analytics(df_clean)

            ai_result, provider_used, fallback_used = await self._generate_insight_payload(
                df_clean=df_clean,
                raw_stats=raw_stats,
                clean_stats=clean_stats,
                cleaning_report=cleaning_report,
                anomaly_detection=anomaly_detection,
                advanced_analytics=advanced_analytics,
                provider=provider,
                custom_prompt=custom_prompt,
            )

            charts_data = self._generate_chart_payload(df_clean, raw_stats, clean_stats)
            result = self._build_frontend_result(
                df_raw=df_raw,
                df_clean=df_clean,
                filename=filename,
                raw_stats=raw_stats,
                clean_stats=clean_stats,
                cleaning_report=cleaning_report,
                ai_result=ai_result,
                charts_data=charts_data,
                anomaly_detection=anomaly_detection,
                advanced_analytics=advanced_analytics,
            )

            await db.update_analysis_status(analysis_id, "completed", result)

            response = {
                "id": analysis_id,
                "filename": filename,
                "status": "completed",
                "result": result,
                "provider_used": provider_used,
            }
            if fallback_used:
                response["fallback_used"] = True
            return response
        except Exception as exc:
            logger.error("[ANALYSIS] failed for %s: %s", filename, exc)
            logger.error(traceback.format_exc())
            await db.update_analysis_status(
                analysis_id,
                "failed",
                {
                    "error": str(exc),
                    "business_insights": "The dataset could not be analyzed because processing failed.",
                },
            )
            raise

    async def get_analysis(self, analysis_id: str) -> Optional[Dict[str, Any]]:
        """Get analysis by ID."""
        return await get_supabase_db().get_analysis(analysis_id)

    async def list_analyses(self, limit: int = 100) -> List[Dict[str, Any]]:
        """List all analyses."""
        return await get_supabase_db().list_analyses(limit)

    async def delete_analysis(self, analysis_id: str) -> bool:
        """Delete a single analysis."""
        return await get_supabase_db().delete_analysis(analysis_id)

    async def delete_all_analyses(self) -> int:
        """Delete all analyses."""
        return await get_supabase_db().delete_all_analyses()

    async def get_portfolio_metrics(self) -> Dict[str, Any]:
        """Summarize stored analyses for the portfolio dashboard."""
        records = await self.list_analyses(limit=500)
        completed = [record for record in records if record.get("status") == "completed"]
        failed = [record for record in records if record.get("status") == "failed"]

        total_rows = 0
        total_columns = 0
        total_charts = 0
        quality_scores: List[float] = []
        advanced_counts = {
            "churn": 0,
            "rfm": 0,
            "predictions": 0,
            "clustering": 0,
        }

        for record in completed:
            result = record.get("analysis_result") or {}
            summary = result.get("summary") or {}
            cleaning = result.get("cleaning_report") or {}
            advanced = result.get("advanced_analytics") or {}

            total_rows += int(summary.get("total_rows") or 0)
            total_columns += int(summary.get("total_columns") or 0)
            total_charts += len(result.get("charts_data") or [])

            quality = cleaning.get("score_after", summary.get("data_quality_score"))
            if isinstance(quality, (int, float)):
                quality_scores.append(float(quality))

            if (advanced.get("churn_analysis") or {}).get("detected"):
                advanced_counts["churn"] += 1
            if (advanced.get("rfm_segmentation") or {}).get("applicable"):
                advanced_counts["rfm"] += 1
            if (advanced.get("predictions") or {}).get("applicable"):
                advanced_counts["predictions"] += 1
            if (advanced.get("clustering") or {}).get("applicable"):
                advanced_counts["clustering"] += 1

        avg_quality = round(sum(quality_scores) / len(quality_scores), 2) if quality_scores else 0

        return {
            "analyses_total": len(records),
            "analyses_completed": len(completed),
            "analyses_failed": len(failed),
            "rows_processed": total_rows,
            "columns_profiled": total_columns,
            "charts_generated": total_charts,
            "avg_quality_score": avg_quality,
            "advanced_signals": advanced_counts,
            "capabilities": [
                "CSV, Excel and JSON ingestion",
                "Automated data cleaning and quality scoring",
                "AI-assisted executive insights with rule-based fallback",
                "Automatic chart recommendation and dashboard generation",
                "Dataset chat, PDF export and PowerPoint export",
                "Local JSON persistence with optional Supabase backend",
            ],
        }

    async def answer_dataset_question(self, dataset_id: str, question: str) -> Dict[str, Any]:
        """Answer common dataset questions using persisted analysis artifacts."""
        record = await self.get_analysis(dataset_id)
        if not record:
            raise ValueError("Dataset not found")

        result = record.get("analysis_result", {})
        clean_stats = result.get("statistics", {}).get("clean", {})
        cleaning_report = result.get("cleaning_report", {})
        advanced = result.get("advanced_analytics", {})
        lower = question.lower()

        if any(term in lower for term in ["average", "avg", "mean", "promedio", "media"]):
            numeric_summary = clean_stats.get("numeric_summary", {})
            ordered = sorted(
                [
                    (col, metrics.get("mean"))
                    for col, metrics in numeric_summary.items()
                    if metrics.get("mean") is not None
                ],
                key=lambda item: abs(item[1]) if item[1] is not None else 0,
                reverse=True,
            )
            if ordered:
                lines = [f"{col}: {value:.2f}" for col, value in ordered[:3]]
                answer = "Top average values in the cleaned dataset:\n- " + "\n- ".join(lines)
            else:
                answer = "I could not find numeric columns with enough data to compute averages."
        elif any(term in lower for term in ["missing", "null", "falt", "vac"]):
            missing = clean_stats.get("missing_values", {})
            columns = missing.get("columns", {})
            top_missing = sorted(columns.items(), key=lambda item: item[1].get("count", 0), reverse=True)
            if top_missing and top_missing[0][1].get("count", 0) > 0:
                lines = [
                    f"{name}: {payload['count']} missing values ({payload['pct']}%)"
                    for name, payload in top_missing[:3]
                ]
                answer = "Remaining missing values after cleaning:\n- " + "\n- ".join(lines)
            else:
                answer = "The cleaned dataset has no remaining missing values in the tracked columns."
        elif any(term in lower for term in ["correlation", "correl", "relacion"]):
            top_corr = self._strongest_correlation(clean_stats.get("correlation_matrix", {}))
            if top_corr:
                left, right, value = top_corr
                answer = (
                    f"The strongest detected relationship is between {left} and {right}, "
                    f"with a correlation of {value:.2f}."
                )
            else:
                answer = "I did not find enough numeric columns to compute a meaningful correlation."
        elif any(term in lower for term in ["quality", "clean", "limpieza", "calidad"]):
            answer = (
                f"Data quality improved from {cleaning_report.get('score_before', 0)}% "
                f"to {cleaning_report.get('score_after', 0)}%. "
                f"The pipeline removed {cleaning_report.get('rows_removed', 0)} rows and applied "
                f"{len(cleaning_report.get('changes_made', []))} cleaning actions."
            )
        elif any(term in lower for term in ["churn", "cancel", "baja"]):
            churn = advanced.get("churn_analysis", {})
            if churn.get("detected"):
                answer = (
                    f"Churn-like behavior was detected with an estimated rate of {churn.get('churn_rate', 0)}%. "
                    + " ".join(churn.get("insights", [])[:2])
                )
            else:
                answer = "The dataset does not expose a reliable churn signal with the available columns."
        elif any(term in lower for term in ["segment", "cluster", "rfm", "persona"]):
            rfm = advanced.get("rfm_segmentation", {})
            clustering = advanced.get("clustering", {})
            if rfm.get("applicable") and rfm.get("segments"):
                segments = ", ".join(
                    f"{seg['name']} ({seg['percentage']}%)" for seg in rfm["segments"][:3]
                )
                answer = f"RFM-style segments detected: {segments}."
            elif clustering.get("applicable") and clustering.get("clusters"):
                segments = ", ".join(
                    f"{seg['name']} ({seg['percentage']}%)" for seg in clustering["clusters"][:3]
                )
                answer = f"Behavioral clusters detected: {segments}."
            else:
                answer = "I do not have enough behavioral or monetary signals to build strong segments."
        elif any(term in lower for term in ["anomaly", "outlier", "anomalia"]):
            anomaly = result.get("anomaly_detection", {})
            answer = (
                f"The pipeline flagged {anomaly.get('detected_rows', 0)} potentially anomalous rows "
                f"({anomaly.get('ratio', 0)}% of the cleaned dataset)."
            )
        else:
            summary = result.get("summary", {})
            insights = result.get("business_insights") or result.get("insights") or ""
            answer = (
                f"This dataset contains {summary.get('total_rows', 0)} cleaned rows across "
                f"{summary.get('total_columns', 0)} columns. "
                f"{insights[:650]}".strip()
            )

        return {"dataset_id": dataset_id, "question": question, "answer": answer}

    async def get_suggested_questions(self, dataset_id: str) -> List[str]:
        """Return dataset-aware question suggestions."""
        record = await self.get_analysis(dataset_id)
        if not record:
            return []

        result = record.get("analysis_result", {})
        clean_stats = result.get("statistics", {}).get("clean", {})
        numeric_columns = clean_stats.get("numeric_columns", [])
        categorical_columns = clean_stats.get("categorical_columns", [])

        suggestions = [
            "What are the main quality issues that were fixed?",
            "Which variables show the strongest relationships?",
            "What should the business prioritize first based on this dataset?",
        ]

        if numeric_columns:
            suggestions.append(f"What is the average behavior of {numeric_columns[0]}?")
        if len(numeric_columns) >= 2:
            suggestions.append(
                f"How do {numeric_columns[0]} and {numeric_columns[1]} move together?"
            )
        if categorical_columns:
            suggestions.append(
                f"Which segments dominate the dataset for {categorical_columns[0]}?"
            )

        return suggestions[:6]

    async def generate_pdf_report(self, dataset_id: str) -> bytes:
        """Generate a PDF report for an analysis."""
        record = await self.get_analysis(dataset_id)
        if not record:
            raise ValueError("Dataset not found")

        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

        result = record.get("analysis_result", {})
        summary = result.get("summary", {})
        cleaning = result.get("cleaning_report", {})
        recommendations = result.get("recommendations", [])
        insights = result.get("business_insights", "")

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, title=record.get("filename", "NexusData Report"))
        styles = getSampleStyleSheet()
        story: List[Any] = [
            Paragraph(f"NexusData Report: {record.get('filename', 'Dataset')}", styles["Title"]),
            Spacer(1, 12),
            Paragraph(
                f"Rows: {summary.get('total_rows', 0)} | Columns: {summary.get('total_columns', 0)}",
                styles["BodyText"],
            ),
            Paragraph(
                f"Data quality improved from {cleaning.get('score_before', 0)}% to {cleaning.get('score_after', 0)}%",
                styles["BodyText"],
            ),
            Spacer(1, 16),
            Paragraph("Executive Insights", styles["Heading2"]),
            Paragraph((insights or "No narrative insights available.").replace("\n", "<br/>"), styles["BodyText"]),
            Spacer(1, 16),
            Paragraph("Recommended Actions", styles["Heading2"]),
        ]

        rec_rows = [["Priority", "Recommendation"]]
        for index, recommendation in enumerate(recommendations[:6], start=1):
            rec_rows.append([str(index), recommendation])

        if len(rec_rows) > 1:
            table = Table(rec_rows, colWidths=[60, 440])
            table.setStyle(
                TableStyle(
                    [
                        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#111827")),
                        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#9ca3af")),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]
                )
            )
            story.append(table)
        else:
            story.append(Paragraph("No recommendations were generated.", styles["BodyText"]))

        doc.build(story)
        return buffer.getvalue()

    async def generate_pptx_report(self, dataset_id: str) -> bytes:
        """Generate a PowerPoint report for an analysis."""
        record = await self.get_analysis(dataset_id)
        if not record:
            raise ValueError("Dataset not found")

        from pptx import Presentation
        from pptx.util import Inches

        result = record.get("analysis_result", {})
        summary = result.get("summary", {})
        cleaning = result.get("cleaning_report", {})
        recommendations = result.get("recommendations", [])
        insights = result.get("business_insights", "")

        prs = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = record.get("filename", "Dataset Report")
        title_slide.placeholders[1].text = "Automated analysis generated by NexusData"

        overview = prs.slides.add_slide(prs.slide_layouts[1])
        overview.shapes.title.text = "Overview"
        overview.placeholders[1].text = (
            f"Rows: {summary.get('total_rows', 0)}\n"
            f"Columns: {summary.get('total_columns', 0)}\n"
            f"Quality score: {cleaning.get('score_after', 0)}%"
        )

        cleaning_slide = prs.slides.add_slide(prs.slide_layouts[1])
        cleaning_slide.shapes.title.text = "Cleaning Impact"
        cleaning_slide.placeholders[1].text = (
            f"Before: {cleaning.get('score_before', 0)}%\n"
            f"After: {cleaning.get('score_after', 0)}%\n"
            f"Rows removed: {cleaning.get('rows_removed', 0)}\n"
            + "\n".join(cleaning.get("changes_made", [])[:4])
        )

        insights_slide = prs.slides.add_slide(prs.slide_layouts[1])
        insights_slide.shapes.title.text = "Key Insights"
        insights_slide.placeholders[1].text = insights[:1200] or "No narrative insights available."

        actions_slide = prs.slides.add_slide(prs.slide_layouts[1])
        actions_slide.shapes.title.text = "Recommended Actions"
        actions_slide.placeholders[1].text = "\n".join(
            f"- {item}" for item in recommendations[:6]
        ) or "No recommendations were generated."

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()

    def get_available_providers(self) -> Dict[str, bool]:
        """Expose the availability of external AI providers."""
        return {"gemini": self.gemini.is_available(), "groq": self.groq.is_available()}

    def normalize_analysis_record(self, record: Dict[str, Any]) -> Dict[str, Any]:
        """Normalize persisted records into the shape the frontend expects."""
        analysis_result = record.get("analysis_result", {}) or {}
        summary = analysis_result.get("summary", {}) if isinstance(analysis_result, dict) else {}
        cleaning = analysis_result.get("cleaning_report", {}) if isinstance(analysis_result, dict) else {}
        dataset_dna = analysis_result.get("dataset_dna", {}) if isinstance(analysis_result, dict) else {}

        if not isinstance(summary, dict):
            summary = {}
        if not isinstance(cleaning, dict):
            cleaning = {}
        if not isinstance(dataset_dna, dict):
            dataset_dna = {}

        if not analysis_result:
            analysis_result = {
                "summary": {
                    "total_rows": 0,
                    "sampled_rows": 0,
                    "total_columns": 0,
                    "data_quality_score": 0,
                    "data_quality_score_raw": 0,
                    "missing_cells": 0,
                    "missing_cells_raw": 0,
                },
                "dataset_dna": {"total_rows": 0, "total_columns": 0},
                "cleaning_report": {
                    "score_before": 0,
                    "score_after": 0,
                    "rows_removed": 0,
                    "original_rows": 0,
                    "final_rows": 0,
                    "improvement": 0,
                    "changes_made": [],
                },
                "business_insights": record.get("error", ""),
            }
        else:
            analysis_result["summary"] = {
                "total_rows": summary.get("total_rows", 0),
                "sampled_rows": summary.get("sampled_rows", summary.get("total_rows", 0)),
                "total_columns": summary.get("total_columns", 0),
                "data_quality_score": summary.get("data_quality_score", cleaning.get("score_after", 0)),
                "data_quality_score_raw": summary.get("data_quality_score_raw", cleaning.get("score_before", 0)),
                "missing_cells": summary.get("missing_cells", 0),
                "missing_cells_raw": summary.get("missing_cells_raw", 0),
            }
            analysis_result["cleaning_report"] = {
                "score_before": cleaning.get("score_before", 0),
                "score_after": cleaning.get("score_after", 0),
                "rows_removed": cleaning.get("rows_removed", 0),
                "original_rows": cleaning.get("original_rows", analysis_result["summary"]["total_rows"]),
                "final_rows": cleaning.get("final_rows", analysis_result["summary"]["total_rows"]),
                "improvement": cleaning.get("improvement", 0),
                "changes_made": cleaning.get("changes_made", []),
            }

        if "dataset_dna" not in analysis_result:
            analysis_result["dataset_dna"] = {
                "total_rows": analysis_result["summary"].get("total_rows", 0),
                "total_columns": analysis_result["summary"].get("total_columns", 0),
                "sampled_rows": analysis_result["summary"].get("sampled_rows", 0),
                "numeric_columns": len(analysis_result.get("statistics", {}).get("clean", {}).get("numeric_columns", [])),
                "categorical_columns": len(analysis_result.get("statistics", {}).get("clean", {}).get("categorical_columns", [])),
                "datetime_columns": len(analysis_result.get("statistics", {}).get("clean", {}).get("datetime_columns", [])),
                "missing_pct": 0,
                "duplicates_removed": analysis_result["cleaning_report"].get("rows_removed", 0),
                "columns_list": analysis_result.get("statistics", {}).get("clean", {}).get("column_names", []),
            }

        return {
            "id": record.get("id"),
            "filename": record.get("filename"),
            "status": record.get("status"),
            "created_at": record.get("created_at"),
            "updated_at": record.get("updated_at"),
            "result": analysis_result,
        }

    async def _generate_insight_payload(
        self,
        df_clean: pd.DataFrame,
        raw_stats: Dict[str, Any],
        clean_stats: Dict[str, Any],
        cleaning_report: Dict[str, Any],
        anomaly_detection: Dict[str, Any],
        advanced_analytics: Dict[str, Any],
        provider: str,
        custom_prompt: Optional[str],
    ) -> Tuple[Dict[str, Any], str, bool]:
        requested = provider if provider in {"gemini", "groq"} else "auto"
        provider_sequence: List[str] = []

        if requested == "gemini":
            provider_sequence = ["gemini", "groq"]
        elif requested == "groq":
            provider_sequence = ["groq", "gemini"]
        else:
            provider_sequence = ["gemini", "groq"]

        for candidate in provider_sequence:
            try:
                if candidate == "gemini" and self.gemini.is_available():
                    result = await self.gemini.analyze_dataset(df_clean, custom_prompt)
                    if not result.get("error"):
                        return result, "gemini", requested != "gemini"
                if candidate == "groq" and self.groq.is_available():
                    result = await self.groq.analyze_dataset(df_clean, custom_prompt)
                    if not result.get("error"):
                        return result, "groq", requested != "groq"
            except Exception as exc:  # pragma: no cover - network/runtime specific
                logger.warning("Provider %s failed, trying next: %s", candidate, exc)

        fallback = self._build_rule_based_insights(
            raw_stats=raw_stats,
            clean_stats=clean_stats,
            cleaning_report=cleaning_report,
            anomaly_detection=anomaly_detection,
            advanced_analytics=advanced_analytics,
        )
        return fallback, "rule-based", True

    def _build_frontend_result(
        self,
        df_raw: pd.DataFrame,
        df_clean: pd.DataFrame,
        filename: str,
        raw_stats: Dict[str, Any],
        clean_stats: Dict[str, Any],
        cleaning_report: Dict[str, Any],
        ai_result: Dict[str, Any],
        charts_data: List[Dict[str, Any]],
        anomaly_detection: Dict[str, Any],
        advanced_analytics: Dict[str, Any],
    ) -> Dict[str, Any]:
        column_types = clean_stats.get("column_types", {})
        missing_total = clean_stats.get("missing_values", {}).get("total", 0)
        raw_missing_total = raw_stats.get("missing_values", {}).get("total", 0)
        dataset_dna = {
            "total_rows": int(len(df_clean)),
            "sampled_rows": int(min(len(df_clean), 5000)),
            "total_columns": int(len(df_clean.columns)),
            "numeric_columns": len(clean_stats.get("numeric_columns", [])),
            "categorical_columns": len(clean_stats.get("categorical_columns", [])),
            "datetime_columns": len(clean_stats.get("datetime_columns", [])),
            "missing_pct": round((missing_total / max(len(df_clean) * max(len(df_clean.columns), 1), 1)) * 100, 2),
            "duplicates_removed": int(df_raw.duplicated().sum()),
            "columns_list": list(df_clean.columns),
        }

        result = {
            "filename": filename,
            "summary": {
                "total_rows": int(len(df_clean)),
                "sampled_rows": int(min(len(df_clean), 5000)),
                "total_columns": int(len(df_clean.columns)),
                "data_quality_score": cleaning_report.get("score_after", 0),
                "data_quality_score_raw": cleaning_report.get("score_before", 0),
                "missing_cells": missing_total,
                "missing_cells_raw": raw_missing_total,
            },
            "dataset_dna": dataset_dna,
            "column_types": column_types,
            "descriptive_statistics": clean_stats.get("numeric_summary", {}),
            "correlation_matrix": clean_stats.get("correlation_matrix", {}),
            "anomaly_detection": anomaly_detection,
            "business_insights": ai_result.get("insights", ""),
            "insights": ai_result.get("insights", ""),
            "summary_text": ai_result.get("summary", ""),
            "sample_data": cleaning_report.get("clean_preview", []),
            "raw_preview": cleaning_report.get("raw_preview", []),
            "clean_preview": cleaning_report.get("clean_preview", []),
            "diff_preview": self._generate_diff_preview(
                cleaning_report.get("raw_preview", []),
                cleaning_report.get("clean_preview", []),
            ),
            "cleaning_report": cleaning_report,
            "charts_data": charts_data,
            "chart_recommendations": [
                {
                    "type": chart.get("type", ""),
                    "x_column": chart.get("x_label", ""),
                    "y_column": chart.get("y_label", ""),
                    "title": chart.get("title", ""),
                    "insight": chart.get("insight", ""),
                }
                for chart in charts_data
            ],
            "recommendations": ai_result.get("recommendations", []),
            "statistics": {"raw": raw_stats, "clean": clean_stats},
            "advanced_analytics": advanced_analytics,
        }
        return self._clean_for_json(result)

    def _normalize_cleaning_report(self, raw_report: Dict[str, Any]) -> Dict[str, Any]:
        """Normalize cleaning report keys for the frontend contract."""
        return {
            "score_before": raw_report.get("quality_before", 0),
            "score_after": raw_report.get("quality_after", 0),
            "rows_removed": raw_report.get("rows_removed", 0),
            "original_rows": raw_report.get("original_rows", 0),
            "final_rows": raw_report.get("final_rows", 0),
            "improvement": raw_report.get("improvement", 0),
            "changes_made": raw_report.get("changes_made", []),
            "steps": raw_report.get("cleaning_steps", []),
            "issues_found": raw_report.get("issues_found", []),
            "quality_before_details": raw_report.get("quality_before_details", {}),
            "raw_preview": raw_report.get("raw_preview", []),
            "clean_preview": raw_report.get("clean_preview", []),
        }

    def _build_dataset_statistics(self, df: pd.DataFrame) -> Dict[str, Any]:
        """Generate descriptive statistics for the UI and export layers."""
        column_types = self._classify_columns(df)
        numeric_columns = [col for col, kind in column_types.items() if kind == "numerical"]
        categorical_columns = [col for col, kind in column_types.items() if kind == "categorical"]
        datetime_columns = [col for col, kind in column_types.items() if kind == "datetime"]

        total_rows = int(len(df))
        total_columns = int(len(df.columns))
        missing_columns = {}
        for col in df.columns:
            count = int(df[col].isna().sum())
            missing_columns[col] = {
                "count": count,
                "pct": round((count / total_rows) * 100, 2) if total_rows else 0,
            }

        numeric_summary = {}
        if numeric_columns:
            describe = df[numeric_columns].describe().to_dict()
            numeric_summary = self._clean_for_json(describe)

        correlation_matrix = {}
        if len(numeric_columns) >= 2:
            correlation_matrix = self._clean_for_json(
                df[numeric_columns].corr(numeric_only=True).round(3).to_dict()
            )

        categorical_summary = {}
        for col in categorical_columns[:6]:
            top_values = df[col].fillna("Unknown").astype(str).value_counts().head(5)
            categorical_summary[col] = [
                {"value": value, "count": int(count)}
                for value, count in top_values.items()
            ]

        return {
            "total_rows": total_rows,
            "total_columns": total_columns,
            "column_names": list(df.columns),
            "column_types": column_types,
            "numeric_columns": numeric_columns,
            "categorical_columns": categorical_columns,
            "datetime_columns": datetime_columns,
            "missing_values": {
                "total": int(df.isna().sum().sum()),
                "columns": missing_columns,
            },
            "numeric_summary": numeric_summary,
            "correlation_matrix": correlation_matrix,
            "categorical_summary": categorical_summary,
        }

    def _classify_columns(self, df: pd.DataFrame) -> Dict[str, str]:
        """Classify columns into business-friendly types."""
        classified: Dict[str, str] = {}
        for col in df.columns:
            series = df[col]
            if is_datetime64_any_dtype(series):
                classified[col] = "datetime"
            elif is_numeric_dtype(series):
                classified[col] = "numerical"
            else:
                sample = series.dropna().astype(str).head(30)
                looks_like_datetime = (
                    not sample.empty
                    and sample.str.contains(r"\d{1,4}[-/]\d{1,2}[-/]\d{1,4}|\d{1,2}:\d{2}", regex=True).mean() >= 0.4
                )
                if looks_like_datetime:
                    converted = pd.to_datetime(series, errors="coerce")
                    valid_ratio = converted.notna().mean() if len(series) else 0
                else:
                    valid_ratio = 0

                if valid_ratio >= 0.6:
                    classified[col] = "datetime"
                else:
                    classified[col] = "categorical"
        return classified

    def _build_rule_based_insights(
        self,
        raw_stats: Dict[str, Any],
        clean_stats: Dict[str, Any],
        cleaning_report: Dict[str, Any],
        anomaly_detection: Dict[str, Any],
        advanced_analytics: Dict[str, Any],
    ) -> Dict[str, Any]:
        """Fallback analyst-style narrative when external AI is unavailable."""
        total_rows = clean_stats.get("total_rows", 0)
        total_columns = clean_stats.get("total_columns", 0)
        score_before = cleaning_report.get("score_before", 0)
        score_after = cleaning_report.get("score_after", 0)
        removed_rows = cleaning_report.get("rows_removed", 0)
        raw_missing = raw_stats.get("missing_values", {}).get("total", 0)
        clean_missing = clean_stats.get("missing_values", {}).get("total", 0)
        top_corr = self._strongest_correlation(clean_stats.get("correlation_matrix", {}))
        dominant_segment = self._dominant_segment(clean_stats.get("categorical_summary", {}))
        prediction = advanced_analytics.get("predictions", {})
        clustering = advanced_analytics.get("clustering", {})

        sections = [
            "## Resumen Ejecutivo",
            (
                f"El dataset limpio quedo con {total_rows} filas y {total_columns} columnas. "
                f"La calidad de datos paso de {score_before}% a {score_after}% "
                f"despues de remover o corregir {removed_rows} filas problematicas."
            ),
            "",
            "## Calidad de Datos",
            (
                f"Antes del procesamiento habia {raw_missing} celdas faltantes; despues de la limpieza "
                f"quedaron {clean_missing}. "
                + (
                    f"Se detectaron {anomaly_detection.get('detected_rows', 0)} filas potencialmente anomalas."
                    if anomaly_detection.get("detected_rows", 0)
                    else "No se detecto un volumen alto de anomalias despues de limpiar."
                )
            ),
            "",
            "## Patrones Relevantes",
        ]

        if top_corr:
            left, right, value = top_corr
            sections.append(
                f"La relacion mas fuerte aparece entre **{left}** y **{right}** con correlacion de {value:.2f}, "
                "lo que sugiere que estos indicadores deben analizarse juntos en cualquier decision operativa."
            )
        else:
            sections.append(
                "No aparecieron correlaciones numericas suficientemente fuertes como para asumir causalidad directa."
            )

        if dominant_segment:
            sections.append(
                f"El segmento mas concentrado es **{dominant_segment['column']} = {dominant_segment['value']}** "
                f"con {dominant_segment['percentage']}% de los registros, lo que puede sesgar conclusiones si se toma como representativo del total."
            )

        sections.extend(["", "## Riesgos y Sesgos Potenciales"])
        bias_notes = self._build_bias_notes(clean_stats, anomaly_detection)
        sections.extend(bias_notes or ["La muestra no muestra un sesgo estructural extremo, pero igual conviene validar decisiones con una ventana temporal mas amplia."])

        sections.extend(["", "## Acciones Recomendadas"])
        recommendations = self._build_recommendations(clean_stats, top_corr, anomaly_detection, advanced_analytics)
        sections.extend([f"- {item}" for item in recommendations])

        if prediction.get("applicable") and prediction.get("r2_score") is not None:
            sections.extend(
                [
                    "",
                    "## Prediccion Automatizada",
                    (
                        f"El modelo exploratorio sobre **{prediction.get('target', 'target')}** obtuvo un R2 aproximado de "
                        f"{prediction.get('r2_score', 0):.2f}, lo que indica una capacidad "
                        + ("util" if prediction.get("r2_score", 0) >= 0.5 else "limitada")
                        + " para explicar la variacion del objetivo."
                    ),
                ]
            )

        if clustering.get("applicable") and clustering.get("clusters"):
            sections.extend(
                [
                    "",
                    "## Segmentacion",
                    "Se detectaron grupos diferenciados en el comportamiento numerico, utiles para priorizar acciones comerciales o de riesgo por cohortes.",
                ]
            )

        insights = "\n".join(sections).strip()
        return {
            "insights": insights,
            "summary": f"Cleaned dataset with {total_rows} rows and a final data quality score of {score_after}%.",
            "recommendations": recommendations,
        }

    def _build_recommendations(
        self,
        clean_stats: Dict[str, Any],
        top_corr: Optional[Tuple[str, str, float]],
        anomaly_detection: Dict[str, Any],
        advanced_analytics: Dict[str, Any],
    ) -> List[str]:
        recommendations: List[str] = []

        if top_corr:
            recommendations.append(
                f"Track {top_corr[0]} and {top_corr[1]} together in reporting because they move in a materially related way."
            )

        if anomaly_detection.get("ratio", 0) >= 5:
            recommendations.append(
                "Review the anomalous rows before using the dataset for forecasting or executive KPI reporting."
            )

        if clean_stats.get("categorical_summary"):
            first_col = next(iter(clean_stats["categorical_summary"]))
            recommendations.append(
                f"Break down KPIs by {first_col} to avoid reading the dataset as a single blended population."
            )

        prediction = advanced_analytics.get("predictions", {})
        if prediction.get("applicable"):
            recommendations.append(
                f"Use {prediction.get('target', 'the detected target')} as a candidate planning KPI and validate the driver variables in the ML section."
            )

        if not recommendations:
            recommendations.append(
                "Keep monitoring data quality on every new upload so operational metrics stay comparable over time."
            )

        return recommendations[:5]

    def _build_bias_notes(
        self,
        clean_stats: Dict[str, Any],
        anomaly_detection: Dict[str, Any],
    ) -> List[str]:
        notes: List[str] = []
        total_rows = clean_stats.get("total_rows", 0)
        if total_rows and total_rows < 50:
            notes.append(
                "La muestra es chica; cualquier patron puede ser inestable y deberia validarse con mas volumen antes de escalar una decision."
            )

        dominant = self._dominant_segment(clean_stats.get("categorical_summary", {}))
        if dominant and dominant["percentage"] >= 70:
            notes.append(
                f"El valor {dominant['value']} domina la columna {dominant['column']} ({dominant['percentage']}%), lo que puede ocultar segmentos minoritarios."
            )

        if anomaly_detection.get("ratio", 0) >= 10:
            notes.append(
                "La proporcion de anomalias es alta; resultados agregados pueden estar inflados por valores extremos."
            )

        return notes

    def _generate_chart_payload(
        self,
        df_clean: pd.DataFrame,
        raw_stats: Dict[str, Any],
        clean_stats: Dict[str, Any],
    ) -> List[Dict[str, Any]]:
        """Generate a diverse chart set for the dashboard."""
        charts: List[Dict[str, Any]] = []
        column_types = clean_stats.get("column_types", {})
        numeric_columns = clean_stats.get("numeric_columns", [])
        categorical_columns = clean_stats.get("categorical_columns", [])
        datetime_columns = clean_stats.get("datetime_columns", [])

        raw_missing_columns = raw_stats.get("missing_values", {}).get("columns", {})
        missing_candidates = [
            (name, payload.get("count", 0))
            for name, payload in raw_missing_columns.items()
            if payload.get("count", 0) > 0
        ]
        if missing_candidates:
            names = [name for name, _ in sorted(missing_candidates, key=lambda item: item[1], reverse=True)[:8]]
            counts = [raw_missing_columns[name]["count"] for name in names]
            charts.append(
                {
                    "type": "bar",
                    "title": "Missing Values by Column",
                    "insight": "Columns with the highest pre-cleaning data quality pressure.",
                    "x": names,
                    "y": counts,
                    "x_label": "Column",
                    "y_label": "Missing values",
                }
            )

        if categorical_columns:
            category_col = categorical_columns[0]
            top_values = (
                df_clean[category_col].fillna("Unknown").astype(str).value_counts().head(8)
            )
            charts.append(
                {
                    "type": "bar",
                    "title": f"Top Categories in {category_col}",
                    "insight": "Shows which segment dominates the cleaned dataset.",
                    "x": top_values.index.tolist(),
                    "y": [int(value) for value in top_values.values.tolist()],
                    "x_label": category_col,
                    "y_label": "Records",
                }
            )

        if numeric_columns:
            numeric_col = numeric_columns[0]
            numeric_values = [
                float(value)
                for value in df_clean[numeric_col].dropna().astype(float).head(500).tolist()
            ]
            if numeric_values:
                charts.append(
                    {
                        "type": "histogram",
                        "title": f"Distribution of {numeric_col}",
                        "insight": "Useful to spot skewness, long tails, and concentration zones.",
                        "values": numeric_values,
                        "x_label": numeric_col,
                        "y_label": "Frequency",
                    }
                )

        if len(numeric_columns) >= 2:
            scatter_df = df_clean[numeric_columns[:2]].dropna().head(300)
            scatter_pairs = scatter_df.values.tolist()
            if scatter_pairs:
                charts.append(
                    {
                        "type": "scatter",
                        "title": f"{numeric_columns[0]} vs {numeric_columns[1]}",
                        "insight": "Compares the two most analysis-ready numeric variables.",
                        "data": scatter_pairs,
                        "x_label": numeric_columns[0],
                        "y_label": numeric_columns[1],
                    }
                )

            corr = clean_stats.get("correlation_matrix", {})
            labels = list(corr.keys())[:6]
            if labels:
                heatmap = [
                    [float(corr[row].get(col, 0)) for col in labels]
                    for row in labels
                ]
                charts.append(
                    {
                        "type": "heatmap",
                        "title": "Correlation Matrix",
                        "insight": "Highlights where numeric drivers move together or in opposite directions.",
                        "data": heatmap,
                        "labels": labels,
                    }
                )

            box_columns = numeric_columns[: min(3, len(numeric_columns))]
            box_data = [
                df_clean[col].dropna().astype(float).head(200).tolist()
                for col in box_columns
                if not df_clean[col].dropna().empty
            ]
            if box_data:
                charts.append(
                    {
                        "type": "boxplot",
                        "title": "Spread of Key Numeric Columns",
                        "insight": "Box plots make outliers and dispersion instantly visible.",
                        "data": box_data,
                        "categories": box_columns[: len(box_data)],
                    }
                )

        if datetime_columns and numeric_columns:
            date_col = datetime_columns[0]
            metric_col = numeric_columns[0]
            series_df = df_clean[[date_col, metric_col]].dropna().copy()
            if not series_df.empty:
                series_df[date_col] = pd.to_datetime(series_df[date_col], errors="coerce")
                series_df = series_df.dropna().sort_values(date_col).head(40)
                grouped = series_df.groupby(series_df[date_col].dt.strftime("%Y-%m-%d"))[metric_col].mean()
                if not grouped.empty:
                    charts.append(
                        {
                            "type": "line",
                            "title": f"{metric_col} over Time",
                            "insight": "Useful for spotting trend acceleration or operational shifts.",
                            "x": grouped.index.tolist(),
                            "y": [float(value) for value in grouped.values.tolist()],
                            "x_label": date_col,
                            "y_label": metric_col,
                        }
                    )
        elif numeric_columns:
            sample_col = numeric_columns[0]
            sample_df = df_clean[[sample_col]].dropna().head(40)
            if not sample_df.empty:
                charts.append(
                    {
                        "type": "area",
                        "title": f"Sequence Trend of {sample_col}",
                        "insight": "A lightweight trend view across the first cleaned records.",
                        "x": [str(index + 1) for index in range(len(sample_df))],
                        "y": [float(value) for value in sample_df[sample_col].astype(float).tolist()],
                        "x_label": "Row index",
                        "y_label": sample_col,
                    }
                )

        type_counts = {"Numerical": 0, "Categorical": 0, "Datetime": 0}
        for kind in column_types.values():
            if kind == "numerical":
                type_counts["Numerical"] += 1
            elif kind == "datetime":
                type_counts["Datetime"] += 1
            else:
                type_counts["Categorical"] += 1
        charts.append(
            {
                "type": "pie",
                "title": "Column Type Mix",
                "insight": "Shows how analysis-ready the schema is by data type.",
                "labels": list(type_counts.keys()),
                "values": list(type_counts.values()),
            }
        )

        return charts[:8]

    def _detect_anomalies(self, df: pd.DataFrame) -> Dict[str, Any]:
        """Detect anomalous rows using Isolation Forest when possible."""
        numeric_columns = [col for col in df.columns if is_numeric_dtype(df[col])]
        if not numeric_columns:
            return {"detected_rows": 0, "ratio": 0}

        numeric_df = df[numeric_columns].dropna()
        if len(numeric_df) < 20:
            return {"detected_rows": 0, "ratio": 0}

        if SKLEARN_AVAILABLE:
            sample = numeric_df.head(2000)
            try:
                model = IsolationForest(contamination="auto", random_state=42)
                predictions = model.fit_predict(sample)
                detected = int((predictions == -1).sum())
                ratio = round((detected / len(sample)) * 100, 2) if len(sample) else 0
                return {"detected_rows": detected, "ratio": ratio}
            except Exception as exc:  # pragma: no cover - ML runtime specifics
                logger.warning("IsolationForest failed, using heuristic anomalies: %s", exc)

        z_scores = np.abs((numeric_df - numeric_df.mean()) / numeric_df.std(ddof=0).replace(0, np.nan))
        detected = int((z_scores > 3).any(axis=1).sum())
        ratio = round((detected / len(numeric_df)) * 100, 2) if len(numeric_df) else 0
        return {"detected_rows": detected, "ratio": ratio}

    def _generate_advanced_analytics(self, df: pd.DataFrame) -> Dict[str, Any]:
        """Generate optional ML and segmentation outputs."""
        return {
            "churn_analysis": self._detect_churn_signal(df),
            "rfm_segmentation": self._build_rfm_segments(df),
            "predictions": self._build_prediction_model(df),
            "clustering": self._build_clusters(df),
        }

    def _detect_churn_signal(self, df: pd.DataFrame) -> Dict[str, Any]:
        candidates = [
            col
            for col in df.columns
            if any(token in col.lower() for token in ["churn", "cancel", "baja", "inactive", "status"])
        ]
        for col in candidates:
            series = df[col].dropna().astype(str).str.lower()
            positive = series.isin({"1", "true", "yes", "churn", "cancelled", "canceled", "inactive", "baja"})
            if positive.any():
                rate = round(float(positive.mean()) * 100, 2)
                return {
                    "detected": True,
                    "churn_rate": rate,
                    "insights": [
                        f"The column {col} behaves like a churn label.",
                        f"Roughly {rate}% of the observed records fall into the churn or inactive state.",
                    ],
                    "risk_level": "high" if rate >= 25 else "medium" if rate >= 10 else "low",
                    "business_impact": "Retention actions should focus on the cohorts driving the inactive share.",
                }
        return {"detected": False, "insights": []}

    def _build_rfm_segments(self, df: pd.DataFrame) -> Dict[str, Any]:
        column_map = {col.lower(): col for col in df.columns}
        recency = next((column_map[key] for key in column_map if "recency" in key or "days_since" in key), None)
        frequency = next((column_map[key] for key in column_map if "frequency" in key or "orders" in key or "purchases" in key), None)
        monetary = next((column_map[key] for key in column_map if "monetary" in key or "revenue" in key or "sales" in key or "amount" in key), None)

        if not (recency and frequency and monetary):
            return {"applicable": False, "insights": []}

        work = df[[recency, frequency, monetary]].dropna().copy()
        if len(work) < 10:
            return {"applicable": False, "insights": []}

        work["r_score"] = pd.qcut(work[recency].rank(method="first"), 3, labels=["At Risk", "Warm", "Recent"])
        work["f_score"] = pd.qcut(work[frequency].rank(method="first"), 3, labels=["Low Frequency", "Medium Frequency", "High Frequency"])
        work["m_score"] = pd.qcut(work[monetary].rank(method="first"), 3, labels=["Low Value", "Mid Value", "High Value"])

        segment_counts = (
            work["m_score"].astype(str).value_counts(normalize=True).mul(100).round(2)
        )
        segments = [
            {"name": name, "count": int((work["m_score"].astype(str) == name).sum()), "percentage": float(pct)}
            for name, pct in segment_counts.items()
        ]

        return {
            "applicable": True,
            "segments": segments[:4],
            "insights": [
                "The dataset exposes enough recency, frequency, and monetary signal for a lightweight RFM cut.",
                "Value-based segmentation can help prioritize retention and upsell motions.",
            ],
            "business_value": "Customer-facing teams can target high-value segments differently from low-frequency cohorts.",
        }

    def _build_prediction_model(self, df: pd.DataFrame) -> Dict[str, Any]:
        numeric_columns = [col for col in df.columns if is_numeric_dtype(df[col])]
        if len(numeric_columns) < 2 or not SKLEARN_AVAILABLE:
            return {"applicable": False, "insights": []}

        target = numeric_columns[-1]
        features = numeric_columns[:-1]
        work = df[features + [target]].dropna().head(3000)
        if len(work) < 30:
            return {"applicable": False, "insights": []}

        try:
            X = work[features]
            y = work[target]
            X_train, X_test, y_train, y_test = train_test_split(
                X, y, test_size=0.2, random_state=42
            )
            model = LinearRegression()
            model.fit(X_train, y_train)
            predictions = model.predict(X_test)
            score = float(r2_score(y_test, predictions))
            top_feature = (
                pd.Series(model.coef_, index=features).abs().sort_values(ascending=False).index[0]
            )
            return {
                "applicable": True,
                "r2_score": round(score, 3),
                "target": target,
                "insights": [
                    f"The strongest modeled driver of {target} is {top_feature}.",
                    "This is an exploratory linear model and should be validated before operational use.",
                ],
                "interpretation": "Higher R2 means the current numeric inputs explain more of the target variation.",
                "recommendation": f"Validate whether {top_feature} is controllable or just correlated with {target}.",
            }
        except Exception as exc:  # pragma: no cover - ML runtime specifics
            logger.warning("Linear regression failed: %s", exc)
            return {"applicable": False, "insights": []}

    def _build_clusters(self, df: pd.DataFrame) -> Dict[str, Any]:
        numeric_columns = [col for col in df.columns if is_numeric_dtype(df[col])]
        if len(numeric_columns) < 2 or not SKLEARN_AVAILABLE:
            return {"applicable": False, "insights": []}

        work = df[numeric_columns].dropna().head(1500)
        if len(work) < 20:
            return {"applicable": False, "insights": []}

        try:
            scaler = StandardScaler()
            scaled = scaler.fit_transform(work)
            n_clusters = 4 if len(work) >= 80 else 3
            model = KMeans(n_clusters=n_clusters, random_state=42, n_init=10)
            labels = model.fit_predict(scaled)
            counts = pd.Series(labels).value_counts(normalize=True).sort_index().mul(100).round(2)
            clusters = [
                {
                    "id": int(cluster_id),
                    "name": f"Cluster {cluster_id + 1}",
                    "size": int((labels == cluster_id).sum()),
                    "percentage": float(percent),
                }
                for cluster_id, percent in counts.items()
            ]
            return {
                "applicable": True,
                "clusters": clusters,
                "insights": [
                    "Numeric behavior separates into distinct operating groups.",
                    "Clusters can be used to prioritize tailored actions instead of one-size-fits-all policies.",
                ],
            }
        except Exception as exc:  # pragma: no cover - ML runtime specifics
            logger.warning("KMeans failed: %s", exc)
            return {"applicable": False, "insights": []}

    def _strongest_correlation(self, matrix: Dict[str, Dict[str, float]]) -> Optional[Tuple[str, str, float]]:
        strongest: Optional[Tuple[str, str, float]] = None
        for left, row in matrix.items():
            for right, value in row.items():
                if left == right:
                    continue
                abs_value = abs(float(value))
                if strongest is None or abs_value > abs(strongest[2]):
                    strongest = (left, right, float(value))
        return strongest

    def _dominant_segment(self, categorical_summary: Dict[str, List[Dict[str, Any]]]) -> Optional[Dict[str, Any]]:
        best: Optional[Dict[str, Any]] = None
        for column, values in categorical_summary.items():
            total = sum(item.get("count", 0) for item in values)
            if total <= 0:
                continue
            top = max(values, key=lambda item: item.get("count", 0))
            percentage = round((top.get("count", 0) / total) * 100, 2)
            candidate = {"column": column, "value": top.get("value", ""), "percentage": percentage}
            if best is None or candidate["percentage"] > best["percentage"]:
                best = candidate
        return best

    def _generate_diff_preview(
        self,
        raw_preview: List[Dict[str, Any]],
        clean_preview: List[Dict[str, Any]],
    ) -> List[Dict[str, Any]]:
        diff: List[Dict[str, Any]] = []
        for row_index in range(min(len(raw_preview), len(clean_preview))):
            raw_row = raw_preview[row_index]
            clean_row = clean_preview[row_index]
            changes = {}
            for key in set(raw_row) | set(clean_row):
                before = raw_row.get(key)
                after = clean_row.get(key)
                if self._normalize_cell(before) != self._normalize_cell(after):
                    changes[key] = {"before": before, "after": after}
            if changes:
                diff.append(
                    {
                        "row": row_index,
                        "changes": changes,
                        "raw": raw_row,
                        "clean": clean_row,
                    }
                )
        return diff

    def _normalize_cell(self, value: Any) -> str:
        if value is None:
            return ""
        return str(value).strip()

    def _clean_for_json(self, obj: Any) -> Any:
        if isinstance(obj, dict):
            return {key: self._clean_for_json(value) for key, value in obj.items()}
        if isinstance(obj, list):
            return [self._clean_for_json(value) for value in obj]
        if isinstance(obj, tuple):
            return [self._clean_for_json(value) for value in obj]
        if isinstance(obj, np.generic):
            return obj.item()
        if isinstance(obj, pd.Timestamp):
            return obj.isoformat()
        if isinstance(obj, float):
            if math.isnan(obj) or math.isinf(obj):
                return None
            return obj
        return obj

    def _parse_dataset(self, file_content: bytes, filename: str) -> pd.DataFrame:
        """Parse dataset from file content."""
        try:
            if filename.endswith(".csv"):
                return pd.read_csv(io.BytesIO(file_content))
            if filename.endswith((".xlsx", ".xls")):
                return pd.read_excel(io.BytesIO(file_content))
            if filename.endswith(".json"):
                return pd.read_json(io.BytesIO(file_content))
            raise ValueError(f"Unsupported file format: {filename}")
        except Exception as exc:
            logger.error("Error parsing dataset %s: %s", filename, exc)
            raise ValueError(f"Failed to parse dataset: {exc}") from exc
